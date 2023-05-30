# -*- coding: UTF-8 -*-
from pptx import Presentation
from pptx.util import Pt

def create_title(prs, title, sub_title):
    # Add a title slide
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    
    title1 = slide1.shapes.title
    title1.text = title
    
    subtitle = slide1.placeholders[1]
    subtitle.text = sub_title

def create_body(prs, slides):
    # Loop through each slide and extract the title and content fields
    for slide in slides:
        slide_title = slide['title']
        slide_content = slide['content']

        # Add a bullet slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(bullet_slide_layout)
        
        title2 = slide2.shapes.title
        title2.text = slide_title
        
        body2 = slide2.shapes.placeholders[1]
        tf = body2.text_frame

        tf.clear()
        
        for content_string in slide_content:
            p = tf.paragraphs[-1]
            
            # set font size
            run = p.add_run()
            run.text = content_string
            font = run.font
            font.name = '微軟正黑體'
            font.size = Pt(18)
            
            p.level = 0
            p = tf.add_paragraph()

if __name__ == '__main__':
    prs = Presentation()
    
    title = '要如何學習機器學習'
    subtitle = 'subtitle'
    create_title(prs, title, subtitle)
    
    slides = []
    slides.append({'title': '進行步驟一', 'content': ['item 1', 'item 2', 'item 3']})
    slides.append({'title': '進行步驟二', 'content': ['item 1', 'item 2', 'item 3']})
    create_body(prs, slides)
    
    prs.save('./1.pptx')    
    
